"""PaddleOCR text extraction for PDF/PPTX/image files."""
import sys
import json
import os
from pathlib import Path

def ocr_pdf(filepath: str) -> str:
    from paddleocr import PPStructure
    from pdf2image import convert_from_path

    engine = PPStructure(lang="ch", use_gpu=False)
    images = convert_from_path(filepath, dpi=200)
    results = []

    for page_num, img in enumerate(images, 1):
        img_path = f"temp_page_{page_num}.png"
        img.save(img_path, "PNG")
        try:
            page_result = engine(img_path)
            texts = []
            for item in page_result:
                if hasattr(item, "text") and item.text:
                    texts.append(item.text)
                elif isinstance(item, dict) and item.get("text"):
                    texts.append(item["text"])
            if texts:
                results.append("\n".join(texts))
        finally:
            if os.path.exists(img_path):
                os.remove(img_path)

    return "\n\n".join(results).strip()


def ocr_pptx(filepath: str) -> str:
    from paddleocr import PPStructure
    from pptx import Presentation
    from io import BytesIO
    import subprocess

    prs = Presentation(filepath)
    engine = PPStructure(lang="ch", use_gpu=False)
    slides_text = []

    # Export slides as images using LibreOffice
    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        outdir = os.path.join(tmpdir, "slides")
        os.makedirs(outdir, exist_ok=True)

        # Convert PPTX to PDF first, then to images
        try:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, filepath],
                timeout=60, check=True, capture_output=True
            )
            pdf_path = os.path.join(tmpdir, Path(filepath).stem + ".pdf")
            if os.path.exists(pdf_path):
                return ocr_pdf(pdf_path)
        except Exception:
            pass

        # Fallback: extract text from PPTX shapes
        texts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                texts.append("\n".join(slide_texts))
        return "\n\n".join(texts).strip()


def ocr_image(filepath: str) -> str:
    from paddleocr import PPStructure
    engine = PPStructure(lang="ch", use_gpu=False)
    result = engine(filepath)
    texts = []
    for item in result:
        if hasattr(item, "text") and item.text:
            texts.append(item.text)
        elif isinstance(item, dict) and item.get("text"):
            texts.append(item["text"])
    return "\n".join(texts).strip()


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)

    filepath = sys.argv[1]
    ext = Path(filepath).suffix.lower()

    try:
        if ext == ".pdf":
            text = ocr_pdf(filepath)
        elif ext in (".pptx", ".ppt"):
            text = ocr_pptx(filepath)
        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
            text = ocr_image(filepath)
        else:
            print(json.dumps({"error": f"Unsupported format: {ext}"}))
            sys.exit(1)

        print(json.dumps({"text": text}, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False))
        sys.exit(1)
