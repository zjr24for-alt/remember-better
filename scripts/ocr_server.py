"""PaddleOCR HTTP server — run with: python scripts/ocr_server.py"""
import json
import os
import sys
import tempfile
from http.server import HTTPServer, BaseHTTPRequestHandler
from pathlib import Path

PORT = 8787


class OCRHandler(BaseHTTPRequestHandler):
    engine = None

    def do_POST(self):
        if self.path != "/ocr":
            self.send_error(404)
            return

        content_length = int(self.headers.get("Content-Length", 0))
        file_bytes = self.rfile.read(content_length)

        # Save uploaded file to temp
        suffix = ".pdf"
        ct = self.headers.get("Content-Type", "")
        if "pdf" in ct:
            suffix = ".pdf"
        elif "pptx" in ct or "powerpoint" in ct:
            suffix = ".pptx"
        elif "png" in ct:
            suffix = ".png"
        elif "jpeg" in ct or "jpg" in ct:
            suffix = ".jpg"

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name

        try:
            text = self.ocr_file(tmp_path, suffix)
            self.send_response(200)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            resp = json.dumps({"text": text}, ensure_ascii=False)
            self.wfile.write(resp.encode("utf-8"))
        except Exception as e:
            self.send_response(500)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(e)}).encode("utf-8"))
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def ocr_file(self, filepath: str, suffix: str) -> str:
        if OCRHandler.engine is None:
            print("[OCR] Loading PaddleOCR engine (first time may take a while)...")
            from paddleocr import PPStructure
            OCRHandler.engine = PPStructure(lang="ch", use_gpu=False, show_log=False)
            print("[OCR] Engine loaded!")

        ext = suffix.lower()

        if ext == ".pdf":
            from pdf2image import convert_from_path
            images = convert_from_path(filepath, dpi=200)
            results = []
            for i, img in enumerate(images):
                img_path = f"{filepath}_page_{i}.png"
                img.save(img_path, "PNG")
                try:
                    items = OCRHandler.engine(img_path)
                    texts = []
                    for item in items:
                        t = item.get("text", "") if isinstance(item, dict) else getattr(item, "text", "")
                        if t:
                            texts.append(t)
                    if texts:
                        results.append("\n".join(texts))
                finally:
                    try:
                        os.unlink(img_path)
                    except Exception:
                        pass
            return "\n\n".join(results)

        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                if slide_texts:
                    texts.append("\n".join(slide_texts))

            # Also try to export slides as images and OCR them
            try:
                import subprocess
                outdir = tempfile.mkdtemp()
                subprocess.run(
                    ["soffice", "--headless", "--convert-to", "pdf", "--outdir", outdir, filepath],
                    timeout=60, check=True, capture_output=True
                )
                pdf_path = os.path.join(outdir, Path(filepath).stem + ".pdf")
                if os.path.exists(pdf_path):
                    pdf_text = self.ocr_file(pdf_path, ".pdf")
                    if pdf_text:
                        texts.append("[OCR from slides]\n" + pdf_text)
            except Exception:
                pass

            return "\n\n".join(texts) if texts else ""

        else:
            # Image
            items = OCRHandler.engine(filepath)
            texts = []
            for item in items:
                t = item.get("text", "") if isinstance(item, dict) else getattr(item, "text", "")
                if t:
                    texts.append(t)
            return "\n".join(texts)


if __name__ == "__main__":
    print(f"[OCR] Starting PaddleOCR server on http://localhost:{PORT}")
    print(f"[OCR] POST files to http://localhost:{PORT}/ocr")
    server = HTTPServer(("0.0.0.0", PORT), OCRHandler)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n[OCR] Server stopped.")
        server.shutdown()
